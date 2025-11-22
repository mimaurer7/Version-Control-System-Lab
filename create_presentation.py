from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
left = top = Inches(1)
width = Inches(8)
height = Inches(2)

# Add title box with styling
title_box = slide1.shapes.add_textbox(left, Inches(2.5), width, height)
title_frame = title_box.text_frame
title_frame.text = "Version Control System Lab"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(54)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)
title_para.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide1.shapes.add_textbox(left, Inches(4), width, Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Professional Git Practices for Secure Development"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(24)
subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
subtitle_para.alignment = PP_ALIGN.CENTER

# Author info
author_box = slide1.shapes.add_textbox(left, Inches(5.5), width, Inches(1))
author_frame = author_box.text_frame
author_frame.text = "Matthew Maurer (mimaurer7)\nCybersecurity Course\nNovember 22, 2024"
author_para = author_frame.paragraphs[0]
author_para.font.size = Pt(16)
author_para.alignment = PP_ALIGN.CENTER

# Slide 2: Overview
slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
title2 = slide2.shapes.title
title2.text = "Lab Overview"
title2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content2 = slide2.placeholders[1]
tf2 = content2.text_frame
tf2.text = "Project Objectives"

points = [
    "Implement Git version control system",
    "Demonstrate Feature Branch Workflow",
    "Integrate security scanning tools",
    "Establish backup and recovery procedures",
    "Document practices for team training"
]

for point in points:
    p = tf2.add_paragraph()
    p.text = point
    p.level = 0
    p.font.size = Pt(20)

# Slide 3: Repository Setup
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title3 = slide3.shapes.title
title3.text = "Setup & Configuration"
title3.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content3 = slide3.placeholders[1]
tf3 = content3.text_frame
tf3.text = "Repository Initialization"

setup_items = [
    "Initialized Git repository (git init)",
    "Configured user credentials",
    "Created comprehensive .gitignore file",
    "Established project directory structure",
    "Made initial commit with all base files"
]

for item in setup_items:
    p = tf3.add_paragraph()
    p.text = item
    p.level = 0
    p.font.size = Pt(18)

# Add note box
note_box = slide3.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
note_frame = note_box.text_frame
note_frame.text = "✓ Repository: https://github.com/mimaurer7/Version-Control-System-Lab"
note_para = note_frame.paragraphs[0]
note_para.font.size = Pt(14)
note_para.font.color.rgb = RGBColor(0, 128, 0)

# Slide 4: Branching Strategy
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title4 = slide4.shapes.title
title4.text = "Feature Branch Workflow"
title4.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content4 = slide4.placeholders[1]
tf4 = content4.text_frame
tf4.text = "Branch Naming Conventions"

branches = [
    "feature/ - New features (feature/add-security-dashboard)",
    "bugfix/ - Bug fixes (bugfix/fix-css-styling)",
    "hotfix/ - Critical production fixes",
    "main - Stable, production-ready code"
]

for branch in branches:
    p = tf4.add_paragraph()
    p.text = branch
    p.level = 0
    p.font.size = Pt(18)

# Add workflow steps
workflow_box = slide4.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
workflow_frame = workflow_box.text_frame
workflow_frame.text = "Workflow: Create branch → Make changes → Commit → Merge to main"
workflow_para = workflow_frame.paragraphs[0]
workflow_para.font.size = Pt(16)
workflow_para.font.italic = True
workflow_para.alignment = PP_ALIGN.CENTER

# Slide 5: Security Practices
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
title5 = slide5.shapes.title
title5.text = "Secure Coding Practices"
title5.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content5 = slide5.placeholders[1]
tf5 = content5.text_frame
tf5.text = "Security Integration"

security = [
    "Pre-commit hooks for automated security scanning",
    "SonarQube - Static code analysis",
    "Snyk - Dependency vulnerability scanning",
    "OWASP Dependency Check - CVE matching",
    "Regular dependency updates (weekly/monthly/immediate)"
]

for item in security:
    p = tf5.add_paragraph()
    p.text = item
    p.level = 0
    p.font.size = Pt(18)

# Slide 6: Backup & Recovery
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
title6 = slide6.shapes.title
title6.text = "Backup & Recovery"
title6.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content6 = slide6.placeholders[1]
tf6 = content6.text_frame
tf6.text = "Remote Backup Strategy"

backup = [
    "GitHub remote repository configured",
    "Automatic backup on every push",
    "Recovery time: 2-5 minutes (full restore)",
    "Tested clone and restore procedures",
    "Multiple recovery scenarios documented"
]

for item in backup:
    p = tf6.add_paragraph()
    p.text = item
    p.level = 0
    p.font.size = Pt(18)

# Slide 7: Training Resources
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
title7 = slide7.shapes.title
title7.text = "Training Resources for Team"
title7.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content7 = slide7.placeholders[1]
tf7 = content7.text_frame
tf7.text = "Essential Git Commands"

commands = [
    "git checkout -b <branch-name> - Create new branch",
    "git add . - Stage all changes",
    "git commit -m 'message' - Commit changes",
    "git push origin <branch> - Push to remote",
    "git merge <branch> - Merge branches"
]

for cmd in commands:
    p = tf7.add_paragraph()
    p.text = cmd
    p.level = 0
    p.font.size = Pt(16)

# Slide 8: Best Practices
slide8 = prs.slides.add_slide(prs.slide_layouts[1])
title8 = slide8.shapes.title
title8.text = "Best Practices & Guidelines"
title8.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content8 = slide8.placeholders[1]
tf8 = content8.text_frame
tf8.text = "Team Guidelines"

practices = [
    "Write clear, descriptive commit messages",
    "Create feature branches for all new work",
    "Never commit secrets or sensitive data",
    "Run security scans before committing",
    "Review code before merging to main",
    "Push changes daily for backup"
]

for practice in practices:
    p = tf8.add_paragraph()
    p.text = practice
    p.level = 0
    p.font.size = Pt(18)

# Slide 9: Summary
slide9 = prs.slides.add_slide(prs.slide_layouts[1])
title9 = slide9.shapes.title
title9.text = "Summary & Key Takeaways"
title9.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

content9 = slide9.placeholders[1]
tf9 = content9.text_frame
tf9.text = "Lab Achievements"

achievements = [
    "✓ Repository initialized and configured",
    "✓ Feature Branch Workflow implemented",
    "✓ Security scanning integrated",
    "✓ Backup and recovery procedures tested",
    "✓ Complete documentation created",
    "✓ Team training resources provided"
]

for achievement in achievements:
    p = tf9.add_paragraph()
    p.text = achievement
    p.level = 0
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 128, 0)

# Save presentation
prs.save('VCS_Lab_Presentation_Matthew_Maurer.pptx')
print("PowerPoint presentation created successfully!")
print("File: VCS_Lab_Presentation_Matthew_Maurer.pptx")